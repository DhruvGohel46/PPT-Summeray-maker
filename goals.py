import io
from pptx import Presentation
from flask import Flask, request, redirect, url_for, flash, render_template, send_file

app = Flask(__name__)
app.secret_key = 'your_secret_key'  # Needed for flashing messages

def create_ppt_with_goal(goal: str) -> io.BytesIO:
    """
    Create a PowerPoint presentation with a title slide that includes the provided goal.
    
    Args:
        goal (str): The presentation goal entered by the user.
    
    Returns:
        io.BytesIO: An in-memory binary stream containing the PPT.
    """
    prs = Presentation()
    # Use the Title Slide layout (typically layout 0)
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    # Set slide title and subtitle
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Presentation"
    subtitle.text = f"Goal: {goal}"
    
    # Save presentation to a BytesIO stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        # Retrieve the presentation goal from the form input (ensure your HTML uses name="goal")
        goal = request.form.get('goal', '').strip()
        if not goal:
            flash("Please enter a goal for the presentation.")
            return redirect(url_for('index'))
        
        # Generate the PPT file using the provided goal
        ppt_file = create_ppt_with_goal(goal)
        
        # Return the generated PPT as a downloadable file
        return send_file(
            ppt_file,
            as_attachment=True,
            attachment_filename="presentation.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    
    # For GET requests, render your form template (which should include an input for 'goal')
    return render_template('index.html')

if __name__ == '__main__':
    app.run(debug=True)
#this will add to app.py file