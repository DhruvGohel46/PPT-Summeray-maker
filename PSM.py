import os
import io
from flask import Flask, request, send_file, render_template, redirect, url_for, flash
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import google.generativeai as genai

from tools.templates import get_available_templates

# Initialize Flask app
app = Flask(__name__)
app.secret_key = 'your_secret_key'  # Replace with a strong secret key
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'output'

# Ensure the directories exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

# List of target audience options for reference
TARGET_AUDIENCE_OPTIONS = {
    'students': "Designed for Students: Engaging and Informative",
    'professionals': "Targeting Professionals: Clear, Concise, Impactful",
    'researchers': "For Researchers: In-depth Analysis and Findings",
    'entrepreneurs': "For Entrepreneurs: Innovative and Future-Focused",
    'general': "General Audience: Broad Overview"
}

# Function to extract text from a PDF
def extract_text_from_pdf(pdf_path):
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

# Function to extract text from a DOCX
def extract_text_from_docx(docx_path):
    doc = docx.Document(docx_path)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

# Summarize text using Gemini 1.5 Flash
def summarize_text(text, goal=None, audience=None):
    api_key = os.environ.get('GOOGLE_API_KEY')  # Retrieve the API key from environment variable
    if not api_key:
        raise ValueError("No API key found. Please set the GOOGLE_API_KEY environment variable.")
        
    genai.configure(api_key=api_key)  # Configure the API with the retrieved key

        
    try:
        genai.configure(api_key=api_key)  # Configure the API with the retrieved key
        model = genai.GenerativeModel("gemini-1.5-flash")
        prompt = f"Please provide a concise summary of the following text:\n\n{text}"
        response = model.generate_content(prompt)
        return response.text if response and response.text else "Summary not generated."
    except Exception as e:
        return f"Error during summarization: {str(e)}"

# Function to add multiple slides to the presentation
def add_multiple_slides(prs, summary_text, num_slides=None):
    words = summary_text.split()
    # Group words into chunks of 6 (each chunk is one line)
    chunked_lines = [' '.join(words[i:i+6]) for i in range(0, len(words), 6)]
    
    if num_slides is None:
        lines_per_slide = 6  # Default: 6 lines per slide
        num_slides = (len(chunked_lines) + lines_per_slide - 1) // lines_per_slide
    else:
        # Calculate approximate lines per slide based on desired slide count
        lines_per_slide = max(1, len(chunked_lines) // num_slides)

    for i in range(num_slides):
        start = i * lines_per_slide
        end = start + lines_per_slide
        slide_lines = chunked_lines[start:end]
        
        slide_layout = prs.slide_layouts[1]  # Title and Content layout.
        slide = prs.slides.add_slide(slide_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

        slide.shapes.title.text = f"Slide {i+1}"
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(20)
        slide.placeholders[1].text = '\n'.join(slide_lines)
        
        # Apply the chosen font to title and content
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
        for paragraph in slide.placeholders[1].text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name



# Create a PowerPoint file with the summary
def create_ppt(summary, output_path, goal=None, audience=None, num_slides=5):

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()  # Set a solid background color
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Set background to white

    # Set title and content
    title_shape = slide.shapes.title  # Set title text
    title_shape.text_frame.text = "Summary"  # Set title text
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Set title font color to black

    title_shape.text_frame.paragraphs[0].font.bold = True  # Make title bold
    title_shape.text_frame.paragraphs[0].font.size = Pt(24)  # Set title font size

    if goal:
        title_shape.text += f" - Goal: {goal}"  # Include the goal in the title
    if len(slide.placeholders) > 1:  # Check if there is a content placeholder

        content_shape = slide.placeholders[1]
        content_shape.text = summary  # Set content text
        content_shape.text_frame.paragraphs[0].font.size = Pt(18)  # Set content font size
        content_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Set content font color to black


    else:
        left = top = Inches(1)
        width = height = Inches(8)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text = summary

    # Call add_multiple_slides to add more slides
    add_multiple_slides(prs, summary, num_slides)  # Add slides based on the calculated number


    prs.save(output_path)

# Create a PowerPoint presentation customized based on the target audience
def create_ppt_with_audience(audience_list: list) -> io.BytesIO:

    prs = Presentation()  # Load the presentation with the selected template

    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()  # Set a solid background color
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Set background to white


    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Presentation"  # Set title text
    title.text_frame.paragraphs[0].font.bold = True  # Make title bold
    title.text_frame.paragraphs[0].font.size = Pt(24)  # Set title font size
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Set title font color to black


    
    # Determine a default subtitle
    custom_subtitle = "Customized Presentation"
    if 'students' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['students']
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Study Tips"
        slide2.placeholders[1].text = "Effective study habits and time management tips."
    elif 'professionals' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['professionals']
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Professional Insights"
        slide2.placeholders[1].text = "Focus on productivity and career growth strategies."
    elif 'researchers' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['researchers']
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Research Findings"
        slide2.placeholders[1].text = "Latest data and detailed analysis."
    elif 'entrepreneurs' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['entrepreneurs']
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Entrepreneurial Strategies"
        slide2.placeholders[1].text = "Innovative ideas and market insights."
    elif 'general' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['general']
    
    subtitle.text = f"Target Audience: {custom_subtitle}"  # Set subtitle text
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)  # Set subtitle font size
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Set subtitle font color to black


    
    # Save presentation to a BytesIO stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

@app.route('/target-audience', methods=['GET', 'POST'])
def target_audience():
    if request.method == 'POST':
        selected_audience = request.form.getlist('audience')
        if not selected_audience:
            flash("Please select at least one target audience.")
            return redirect(url_for('target_audience'))
        
        ppt_file = create_ppt_with_audience(selected_audience)
        return send_file(
            ppt_file,
            as_attachment=True,
            attachment_filename="presentation.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    
    return render_template('target_audience.html', options=TARGET_AUDIENCE_OPTIONS)

# Route for home page
@app.route('/', methods=['GET', 'POST'])
def index():
    if request.method == 'POST':
        if 'datafile' not in request.files or file.filename == '':

            return "No file part", 400
        file = request.files['datafile']
        if file.filename == '':
            return "No selected file", 400
        
        filename = secure_filename(file.filename)
        upload_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(upload_path)

        presentation_goal = request.form.get('goal')  # Retrieve the goal input from the form

        # Determine file type        
        ext = os.path.splitext(filename)[1].lower()
        if ext == '.pdf' or ext == '.docx':

            text = extract_text_from_pdf(upload_path)
        elif ext == '.docx':
            text = extract_text_from_docx(upload_path)
        else:
            return "Unsupported file type. Please upload a PDF or DOCX file.", 400

        if text:
            summary = summarize_text(text, goal=presentation_goal, audience=request.form.getlist('audience'))  # Pass goal and audience to summarize_text
        else:
            return "Error processing the file.", 400

        ppt_path = os.path.join(app.config['OUTPUT_FOLDER'], 'summary_presentation.pptx')        
        create_ppt(summary, ppt_path, presentation_goal, request.form.getlist('audience'))  # Pass the presentation goal and audience to create_ppt

        
        return send_file(ppt_path, as_attachment=True)  # Send the generated PPT back to the user

    templates = get_available_templates()
    return render_template('index.html', templates=templates)  # Pass templates to the index.html
    
    # Ensure that the templates exist
    if not os.path.exists('templates/target_audience.html') or not os.path.exists('templates/index.html'):
        logger.error("Required template files are missing.")


if __name__ == '__main__':
    app.run(host='127.0.0.1', debug=True)
