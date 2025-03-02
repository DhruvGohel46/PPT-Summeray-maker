import io
from flask import Flask, request, render_template, redirect, url_for, flash, send_file
from pptx import Presentation

app = Flask(__name__)
app.secret_key = 'your_secret_key'  # Replace with a strong secret key

# List of target audience options for reference
TARGET_AUDIENCE_OPTIONS = {
    'students': "Designed for Students: Engaging and Informative",
    'professionals': "Targeting Professionals: Clear, Concise, Impactful",
    'researchers': "For Researchers: In-depth Analysis and Findings",
    'entrepreneurs': "For Entrepreneurs: Innovative and Future-Focused",
    'general': "General Audience: Broad Overview"
}

def create_ppt_with_audience(audience_list: list) -> io.BytesIO:
    """
    Create a PowerPoint presentation customized based on the target audience.
    
    Args:
        audience_list (list): A list of selected target audience keys.
    
    Returns:
        io.BytesIO: An in-memory binary stream containing the PPT.
    """
    prs = Presentation()
    # Create the title slide (layout 0)
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Presentation"
    
    # Determine a default subtitle
    custom_subtitle = "Customized Presentation"
    # Check for specific audience values
    if 'students' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['students']
        # Add an extra slide for Students
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Study Tips"
        slide2.placeholders[1].text = "Effective study habits and time management tips."
    elif 'professionals' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['professionals']
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Professional Insights"
        slide2.placeholders[1].text = "Focus on productivity and career growth strategies."
    elif 'researchers' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['researchers']
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Research Findings"
        slide2.placeholders[1].text = "Latest data and detailed analysis."
    elif 'entrepreneurs' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['entrepreneurs']
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Entrepreneurial Strategies"
        slide2.placeholders[1].text = "Innovative ideas and market insights."
    elif 'general' in audience_list:
        custom_subtitle = TARGET_AUDIENCE_OPTIONS['general']
    else:
        custom_subtitle = "Customized Presentation"
    
    subtitle.text = f"Target Audience: {custom_subtitle}"
    
    # Save presentation to a BytesIO stream
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream

@app.route('/target-audience', methods=['GET', 'POST'])
def target_audience():
    if request.method == 'POST':
        # Retrieve the list of selected target audience options (via checkboxes with name "audience")
        selected_audience = request.form.getlist('audience')
        if not selected_audience:
            flash("Please select at least one target audience.")
            return redirect(url_for('target_audience'))
        
        # Generate the PPT using the selected target audience options
        ppt_file = create_ppt_with_audience(selected_audience)
        return send_file(
            ppt_file,
            as_attachment=True,
            attachment_filename="presentation.pptx",
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    
    # For GET, render a template that includes the checkboxes for target audience selection
    return render_template('target_audience.html', options=TARGET_AUDIENCE_OPTIONS)

if __name__ == '__main__':
    app.run(debug=True)
#this will add to app.py file
