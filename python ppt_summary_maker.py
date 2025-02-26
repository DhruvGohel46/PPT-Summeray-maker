# PPT Summary Maker

import sys
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from pptx.util import Inches
from transformers import BartTokenizer, BartForConditionalGeneration
import torch

def extract_text_from_pdf(pdf_path):
    reader = PdfReader(pdf_path)
    text = ""
    for page in reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text

def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
    except Exception as e:
        print(f"Error opening the document: {e}")
        sys.exit(1)
    text = ""
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def summarize_text(text):
    # Use facebook/bart-large-cnn model for summarization
    tokenizer = BartTokenizer.from_pretrained("facebook/bart-large-cnn")
    model = BartForConditionalGeneration.from_pretrained("facebook/bart-large-cnn")
    inputs = tokenizer(text, return_tensors="pt", max_length=1024, truncation=True)
    summary_ids = model.generate(
        inputs.input_ids, 
        num_beams=4, 
        max_length=150, 
        min_length=50, 
        early_stopping=True
    )
    summary = tokenizer.decode(summary_ids[0], skip_special_tokens=True)
    return summary

def create_ppt(summary, output_path='summary_presentation.pptx'):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = "Summary"
    
    if len(slide.placeholders) > 1:
        content_shape = slide.placeholders[1]
        content_shape.text = summary
    else:
        left = top = Inches(1)
        width = height = Inches(8)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text = summary

    prs.save(output_path)
    print(f"Presentation saved as {output_path}")

def main():
    file_type = input("Enter file type (pdf/docx): ").strip().lower()
    file_path = input("Enter file path: ").strip()
    
    if file_type == 'pdf':
        text = extract_text_from_pdf(file_path)
    elif file_type == 'docx':
        text = extract_text_from_docx(file_path)
    else:
        print("Unsupported file type.")
        sys.exit(1)
    
    summary = summarize_text(text)
    create_ppt(summary)

if __name__ == '__main__':
    main()
