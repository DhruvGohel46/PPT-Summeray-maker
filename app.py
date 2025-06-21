import os
import io 
import logging
from flask import Flask, request, send_file, render_template, redirect, url_for, flash, jsonify, session
from werkzeug.utils import secure_filename
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import google.generativeai as genai

# Initialize Flask app
app = Flask(__name__)

# Set up logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app.secret_key = 'your_secret_key'  # Replace with a strong secret key
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'output'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Ensure the directories exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

# List of target audience options
TARGET_AUDIENCE_OPTIONS = {
    'students': "Designed for Students: Engaging and Informative",
    'professionals': "Targeting Professionals: Clear, Concise, Impactful",
    'researchers': "For Researchers: In-depth Analysis and Findings",
    'entrepreneurs': "For Entrepreneurs: Innovative and Future-Focused",
    'general': "General Audience: Broad Overview"
}

# List of font options
FONT_OPTIONS = {
    # Default options (shown first)
    'default': [
        'Calibri', 'Arial', 'Times New Roman', 'Verdana', 
        'Tahoma', 'Georgia', 'Segoe UI', 'Cambria',
        'Century Gothic', 'Garamond'
    ],
    # Additional options (shown when "More options" is clicked)
    'more': [
        'Palatino Linotype', 'Book Antiqua', 'Trebuchet MS', 'Courier New',
        'Franklin Gothic Medium', 'Lucida Sans', 'Constantia', 'Comic Sans MS',
        'Candara', 'MS Sans Serif', 'Arial Black', 'Impact', 'Rockwell',
        'Bookman Old Style', 'Copperplate Gothic', 'Bahnschrift', 'Consolas',
        'Lucida Console', 'Baskerville Old Face', 'Bookshelf Symbol 7'
    ]
}

# Function to extract text from a PDF
def extract_text_from_pdf(pdf_path):
    try:
        reader = PdfReader(pdf_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text
    except Exception as e:
        logger.error(f"Failed to read PDF file: {str(e)}")
        logger.info("Request method: %s", request.method)
        logger.info("Request data: %s", request.form)
        return None

# Function to extract text from a DOCX
def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        logger.error(f"Failed to read DOCX file: {str(e)}")
        logger.info("Request method: %s", request.method)
        logger.info("Request data: %s", request.form)
        return None

# Summarize text using Google's Gemini model
def summarize_text(text, goal=None, audience=None, num_slides=None):
    try:
        # Get API key from environment variable
        api_key = os.environ.get('GEMINI_API_KEY')
        if not api_key:
            logger.warning("No API key found in environment. Using fallback method.")
            # Note: In production, remove this hardcoded key and use only environment variables
            api_key = "AIzaSyC8YQ0xsE1RfPByU6NeVZZrWItXpSgeHpw"
        
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel("gemini-2.0-flash")

        # Ensure goal and audience are not None or empty
        goal_text = goal if goal else "General goal"
        audience_text = audience if audience else "General audience"
        slides_count = f" Create exactly {num_slides} slides." if num_slides else " Create around 8-15 slides."

        # Updated prompt to ensure proper formatting with concise bullet points
        prompt = f"""Please summarize the following text into a structured PowerPoint presentation format.

Requirements:
- The goal of this presentation is: {goal_text}
- The target audience is: {audience_text}{slides_count}
- IMPORTANT: Each slide MUST start with 'Slide X Title: [Title]' where X is the slide number
- IMPORTANT: Each slide MUST have EXACTLY 5 bullet points - no more, no less
- IMPORTANT: Each bullet point MUST start with a dash '-'
- CRUCIAL: Keep each bullet point VERY concise - ideally 60 characters or less not less than 50, maximum 80 characters
- Use short phrases instead of complete sentences where possible
- Avoid long, wordy explanations in bullet points
- Prefer keywords and key phrases that convey the essential information
- Do not include any text that is not part of a slide title or bullet point
- And try to Cover Maximum points 
- if title or subtitle already given in file so consider that important for summarizing
- stricky follow all the rules 
- make the best summery of the text and not to miss any part of given text


Example format (exactly follow this pattern):
Slide 1 Title: Introduction
- Main topic: [concise description]
- Key audience benefit
- Core problem addressed
- Solution overview
- Expected outcomes

Slide 2 Title: Key Findings
- Finding 1: [brief result]
- Finding 2: [brief result]
- Primary data point: [specific number/stat]
- Secondary observation
- Implications for stakeholders

Text to summarize:
{text}"""
        
        logger.info("Sending prompt to Gemini API")
        response = model.generate_content(prompt)
        summary = response.text if response and response.text else "Summary not generated."
        logger.info(f"Received summary from Gemini API (first 100 chars): {summary[:100]}...")
        return summary
    except Exception as e:
        logger.error(f"Error during summarization: {str(e)}")
        logger.info("Request method: %s", request.method)
        logger.info("Request data: %s", request.form)
        return f"Error during summarization: {str(e)}"
# Function to estimate the number of slides based on the summary
def estimate_slides(summary):

    # Count the number of slide titles in the structured summary
    lines = summary.strip().split('\n')
    slide_count = 0
    
    for line in lines:
        if line.startswith('Slide') and 'Title:' in line:
            slide_count += 1
    
    logger.info(f"Estimated {slide_count} slides from summary")
    # Return at least 1 slide
    return max(1, slide_count)

# Parse the structured summary into a list of slides with titles and bullet points
def parse_structured_summary(summary):
    slides = []
    current_slide = None
    
    # Log summary for debugging
    logger.info(f"Parsing summary with length: {len(summary)}")
    
    # Split the summary into lines
    lines = summary.strip().split('\n')
    logger.info(f"Found {len(lines)} lines in summary")
    
    for line_num, line in enumerate(lines):
        line = line.strip()
        if not line:
            continue
        
        # Log each line for debugging
        logger.info(f"Processing line {line_num}: {line[:50]}...")
        
        # Check if this line is a slide title
        if line.startswith('Slide') and 'Title:' in line:
            logger.info(f"Found slide title: {line}")
            # If we were already processing a slide, add it to our list
            if current_slide:
                slides.append(current_slide)
                logger.info(f"Added slide: {current_slide['title']} with {len(current_slide['bullets'])} bullets")
            
            # Extract the title text after the "Slide X Title:" prefix
            title_parts = line.split('Title:', 1)
            title = title_parts[1].strip() if len(title_parts) > 1 else "Slide"
            
            # Start a new slide
            current_slide = {
                'title': title,
                'bullets': []
            }
        
        # Check if this line is a bullet point
        elif line.startswith('-') and current_slide:
            # Add this bullet point to the current slide
            bullet_text = line[1:].strip()  # Remove the dash and trim whitespace
            if bullet_text:  # Only add non-empty bullet points
                current_slide['bullets'].append(bullet_text)
                logger.info(f"Added bullet: {bullet_text[:30]}...")
        
        # Special handling for lines that should be bullets but don't start with dash
        elif current_slide and not line.startswith('Slide'):
            # This might be a continuation of a bullet point or a malformatted bullet
            logger.info(f"Found text that might be a bullet: {line[:30]}...")
            # If it starts with a bullet-like character, treat as a new bullet
            if line.startswith('•') or line.startswith('*'):
                bullet_text = line[1:].strip()
                if bullet_text:
                    current_slide['bullets'].append(bullet_text)
                    logger.info(f"Added alternative bullet: {bullet_text[:30]}...")
    
    # Don't forget to add the last slide
    if current_slide:
        slides.append(current_slide)
        logger.info(f"Added final slide: {current_slide['title']} with {len(current_slide['bullets'])} bullets")
    
    logger.info(f"Total slides parsed: {len(slides)}")
    return slides

# Create a PowerPoint file with the structured summary
def create_ppt(summary, output_path, title=None, goal=None, audience=None, font_name='Calibri', num_slides=None):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    title_only = prs.slides.add_slide(title_slide_layout)
    title_placeholder = title_only.shapes.title
    
    # Set title
    title_text = title if title else "Presentation Summary"
    title_placeholder.text = title_text
    
    # Remove subtitle box by setting it to empty string if it exists
    for shape in title_only.shapes:
        if shape.has_text_frame and shape != title_placeholder:
            shape.text = ""  # Set to empty instead of trying to remove
    
    # Apply font to title
    if hasattr(title_placeholder, "text_frame"):
        for paragraph in title_placeholder.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(50)  # Set title slide Font Size to 50 
                run.font.bold = True  # Make Title Slide Text Bold 
    
    # Parse the structured summary and create slides
    slides_content = parse_structured_summary(summary)
    logger.info(f"Creating PowerPoint with {len(slides_content)} slides")
    
    # Analyze all content to determine optimal font size that works for all slides
    # First, gather all bullet points
    all_bullet_points = []
    for slide_content in slides_content:
        bullet_points = slide_content.get('bullets', [])
        all_bullet_points.extend(bullet_points)
    
    # Determine the optimal font size based on the longest text in any slide
    max_line_length = max([len(point) for point in all_bullet_points if point]) if all_bullet_points else 0
    
    # Calculate content font size - adjust these thresholds based on testing
    # These values ensure text fits within standard slide dimensions
    if max_line_length > 80:
        content_font_size = Pt(20)
    elif max_line_length > 75:
        content_font_size = Pt(21)
    elif max_line_length > 70:
        content_font_size = Pt(22)
    elif max_line_length > 65:
        content_font_size = Pt(23)
    elif max_line_length > 60:
        content_font_size = Pt(24)
    elif max_line_length > 55:
        content_font_size = Pt(25)
    elif max_line_length > 50:
        content_font_size = Pt(28)
    else:
        content_font_size = Pt(30)
    
    logger.info(f"Using consistent content font size of {content_font_size.pt} points for all slides")
    
    # Create content slides from the parsed structure
    for slide_content in slides_content:
        slide_title = slide_content.get('title', 'Slide')
        bullet_points = slide_content.get('bullets', [])
        
        # Skip empty slides
        if not slide_title and not bullet_points:
            logger.info("Skipping empty slide")
            continue
        
        # Add a new slide with title and content layout (layout index 1)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        logger.info(f"Added slide with title: {slide_title}")
        
        # Set the slide title
        if slide.shapes.title:
            title_shape = slide.shapes.title
            slide.shapes.title.text = slide_title
            
            # Apply formatting to title text: bold and 40pt
            for paragraph in title_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(38)
                    run.font.bold = True
        
        # Ensure exactly 5 bullet points
        if len(bullet_points) < 5:
            # Add empty bullets to reach 5
            while len(bullet_points) < 5:
                bullet_points.append("")
        elif len(bullet_points) > 5:
            # Keep only the first 5 bullets
            bullet_points = bullet_points[:5]
        
        # Add bullet points to the content placeholder
        if len(slide.placeholders) > 1:
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()  # Clear any existing text
            
            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = point
                p.level = 0  # Set to first level bullet
                
                # Apply consistent font formatting across all slides
                for run in p.runs:
                    run.font.name = font_name
                    run.font.size = content_font_size  # Use the consistent font size
                
                logger.info(f"Added bullet point {i+1}: {point[:30]}...")
    
    try:
        prs.save(output_path)
        logger.info(f"PowerPoint saved to {output_path}.")
        return True
    except Exception as e:
        logger.error(f"Error saving PowerPoint: {str(e)}")
        return False
    
@app.route('/', methods=['GET'])
def index():
    """Render the initial page with file type selection"""
    return render_template('index.html')

@app.route('/setup', methods=['POST'])
def setup():
    """Handle the file type selection and render the form"""
    file_type = request.form.get('file_type')
    if file_type not in ['pdf', 'docx']:
        flash("Please select a valid file type.")
        return redirect(url_for('index'))
    
    session['file_type'] = file_type
    return render_template(
        'setup_form.html', 
        file_type=file_type,
        audience_options=TARGET_AUDIENCE_OPTIONS,
        font_options=FONT_OPTIONS
    )

def estimate_slides(summary):
    # Count the number of slide titles in the structured summary
    lines = summary.strip().split('\n')
    slide_count = 1
    
    for line in lines:
        if line.startswith('Slide') and 'Title:' in line:
            slide_count += 1
    
    logger.info(f"Estimated {slide_count} slides from summary")
    # Return at least 1 slide
    return max(1, slide_count)

@app.route('/process', methods=['POST'])
def process():
    """Process the uploaded file and form data"""
    # Get form data
    title = request.form.get('title', 'Presentation')
    audience = request.form.get('audience', 'general')
    goal = request.form.get('goal', '')
    font_name = request.form.get('font', 'Calibri')
    file_type = session.get('file_type', 'pdf')
    
    # Check if file was uploaded
    if 'file' not in request.files or request.files['file'].filename == '':
        flash("Please upload a file.")
        return redirect(url_for('setup'))
    
    file = request.files['file']
    filename = secure_filename(file.filename)
    
    # Check file extension
    ext = os.path.splitext(filename)[1].lower()
    expected_ext = '.pdf' if file_type == 'pdf' else '.docx'
    
    if ext != expected_ext:
        flash(f"Please upload a {file_type.upper()} file. You selected {file_type} but uploaded a file with extension {ext}.")
        return redirect(url_for('setup'))
    
    # Save the file
    upload_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
    file.save(upload_path)
    
    # Extract text from the file
    if file_type == 'pdf':
        text = extract_text_from_pdf(upload_path)
    else:  # docx
        text = extract_text_from_docx(upload_path)
    
    if not text:
        flash(f"Error: Could not extract text from the {file_type.upper()} file.")
        return redirect(url_for('setup'))
    
    # Summarize the text
    summary = summarize_text(text, goal=goal, audience=audience)
    
    # Estimate number of slides
    estimated_slides = estimate_slides(summary)
    
    # Save data in session for later use
    session['summary'] = summary
    session['title'] = title
    session['audience'] = audience
    session['goal'] = goal
    session['font_name'] = font_name
    session['estimated_slides'] = estimated_slides
    
    # Show confirmation page
    return render_template('confirm.html', estimated_slides=estimated_slides)

@app.route('/confirm', methods=['POST'])
def confirm():
    """Handle slide count confirmation and generate the final PPT"""
    response = request.form.get('response', '').lower()
    
    # Get saved data from session
    summary = session.get('summary', '')
    title = session.get('title', 'Presentation')
    audience = session.get('audience', 'general')
    goal = session.get('goal', '')
    font_name = session.get('font_name', 'Calibri')
    estimated_slides = session.get('estimated_slides', 5)
    
    if not summary:
        flash("Session expired. Please upload your file again.")
        return redirect(url_for('index'))
    
    num_slides = None  # Default to auto-detect
    
    if response == 'no':
        # User wants custom slide count
        custom_slides = request.form.get('custom_slides', '')
        try:
            num_slides = int(custom_slides)
            if num_slides <= 0:
                raise ValueError("Slide count must be positive")
            # Regenerate summary with the desired number of slides
            logger.info(f"Regenerating summary for {num_slides} slides")
            text = session.get('original_text', '')
            if text:
                summary = summarize_text(text, goal=goal, audience=audience, num_slides=num_slides)
                session['summary'] = summary
            else:
                logger.warning("Original text not found in session, using existing summary")
        except ValueError:
            flash("Please enter a valid number of slides.")
            return render_template('confirm.html', estimated_slides=estimated_slides)
    
    logger.info("Generating the PowerPoint presentation.")
    # Generate the PPT
    output_filename = f"{secure_filename(title)}.pptx"
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    
    success = create_ppt(
        summary, 
        output_path, 
        title=title, 
        goal=goal, 
        audience=audience, 
        font_name=font_name,
        num_slides=num_slides
    )
    
    if not success:
        logger.error("Failed to generate the PowerPoint presentation.")
        flash("Error generating the presentation. Please try again.")
        return redirect(url_for('index'))
    
    # Check if the output file exists before sending
    if not os.path.exists(output_path):
        flash("Error: The PowerPoint file was not generated.")
        return redirect(url_for('index'))

    # Clear session data only after successful download
    session.clear()  
    
    # Send the file
    logger.info(f"Sending the PowerPoint file: {output_filename}")
    return send_file(
        output_path,
        as_attachment=True,
        download_name=output_filename,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

@app.errorhandler(413)
def request_entity_too_large(error):
    """Handle file too large error"""
    flash("The file you uploaded is too large. Please upload a file smaller than 16MB.")
    return redirect(url_for('index'))

if __name__ == '__main__':
    app.run(host='127.0.0.1', debug=True)